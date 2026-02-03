"""Worker xử lý export jobs với ProcessPool."""

import asyncio
import os
from concurrent.futures import ProcessPoolExecutor
from datetime import datetime
from pathlib import Path
from typing import Any

from app.core.config import EXPORTS_DIR
from app.db.session import async_session_maker, init_db
from app.services import export_job as export_job_service

DIAGRAM_EXPORT_TYPES = {"l1_diagram", "l2_diagram", "l3_diagram"}
FILE_EXPORT_TYPES = {"device_file", "master_file"}


def _generate_pptx(file_path: Path, title: str, metadata: dict[str, Any]) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    from app.services.link_palette import get_link_palette_rgb

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    shapes = slide.shapes
    if shapes.title:
        shapes.title.text = title
    lines = [f"{k}: {v}" for k, v in metadata.items()]
    if len(shapes.placeholders) > 1:
        body = shapes.placeholders[1]
        body.text = "\n".join(lines)
    else:
        textbox = shapes.add_textbox(0, 0, presentation.slide_width, presentation.slide_height)
        textbox.text = "\n".join(lines)
    palette = get_link_palette_rgb()
    legend_left = Inches(0.6)
    legend_top = Inches(1.7)
    swatch_size = Inches(0.22)
    row_height = Inches(0.28)

    for idx, (purpose, rgb) in enumerate(palette.items()):
        top = legend_top + row_height * idx
        swatch = shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left, top, swatch_size, swatch_size)
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = RGBColor(*rgb)
        swatch.line.color.rgb = RGBColor(*rgb)

        label_box = shapes.add_textbox(legend_left + Inches(0.32), top - Inches(0.02), Inches(2.8), swatch_size)
        text_frame = label_box.text_frame
        text_frame.clear()
        run = text_frame.paragraphs[0].add_run()
        run.text = purpose
        run.font.size = Pt(10)

    presentation.save(str(file_path))


def _generate_xlsx(file_path: Path, title: str, metadata: dict[str, Any]) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "export"
    sheet.append(["key", "value"])
    sheet.append(["title", title])
    for key, value in metadata.items():
        sheet.append([key, str(value)])
    workbook.save(str(file_path))


def generate_export_file(
    export_type: str,
    project_id: str,
    options: dict[str, Any],
    exports_dir: str,
) -> tuple[str, str, int]:
    """Tạo file export (PPTX/Excel) trong ProcessPool."""
    export_root = Path(exports_dir)
    export_root.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    format_name = options.get("format")

    if export_type in DIAGRAM_EXPORT_TYPES:
        format_name = format_name or "pptx"
        file_name = f"{export_type}_{timestamp}.{format_name}"
        file_path = export_root / file_name
        metadata = {
            "project_id": project_id,
            "export_type": export_type,
            "mode": options.get("mode"),
            "theme": options.get("theme"),
            "format": format_name,
        }
        _generate_pptx(file_path, "BSV Network Sketcher Export", metadata)
    elif export_type in FILE_EXPORT_TYPES:
        format_name = format_name or "xlsx"
        file_name = f"{export_type}_{timestamp}.{format_name}"
        file_path = export_root / file_name
        metadata = {
            "project_id": project_id,
            "export_type": export_type,
            "format": format_name,
        }
        _generate_xlsx(file_path, "BSV Network Sketcher Export", metadata)
    else:
        raise ValueError("Export type không hợp lệ")

    return str(file_path), file_name, file_path.stat().st_size


async def _process_job(executor: ProcessPoolExecutor) -> bool:
    """Xử lý 1 job pending. Trả về True nếu có job."""
    async with async_session_maker() as db:
        job = await export_job_service.get_next_pending_job(db)
        if not job:
            return False

        job_id = job.id
        project_id = job.project_id
        export_type = job.export_type
        options = export_job_service.parse_options(job.options_json) or {}
        await export_job_service.mark_processing(db, job)

    loop = asyncio.get_running_loop()
    try:
        file_path, file_name, file_size = await loop.run_in_executor(
            executor,
            generate_export_file,
            export_type,
            project_id,
            options,
            EXPORTS_DIR,
        )
        async with async_session_maker() as db:
            job = await export_job_service.get_job(db, job_id)
            if job:
                await export_job_service.mark_completed(
                    db,
                    job,
                    file_path=file_path,
                    file_name=file_name,
                    file_size=file_size,
                    message="Export hoàn tất",
                )
    except Exception as exc:
        async with async_session_maker() as db:
            job = await export_job_service.get_job(db, job_id)
            if job:
                await export_job_service.mark_failed(db, job, error_message=str(exc))
    return True


async def run_worker() -> None:
    """Background worker để xử lý export jobs (PPTX/Excel)."""
    poll_interval = float(os.getenv("EXPORT_POLL_INTERVAL", "2"))
    max_workers = int(os.getenv("EXPORT_MAX_WORKERS", "2"))
    await init_db()

    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        while True:
            handled = await _process_job(executor)
            if not handled:
                await asyncio.sleep(poll_interval)


if __name__ == "__main__":
    asyncio.run(run_worker())
