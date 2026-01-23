#!/usr/bin/env python3
"""
Trích xuất metadata từ PPTX/Excel để so sánh regression.

Sử dụng:
    python scripts/extract_metadata.py golden-files/outputs/small/l1_diagram.pptx
    python scripts/extract_metadata.py golden-files/outputs/small/master_file.xlsx

Yêu cầu:
    pip install python-pptx openpyxl

Output: JSON metadata để so sánh tự động
"""
import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


def extract_pptx_metadata(pptx_path: str) -> dict:
    """
    Trích xuất metadata từ PPTX.

    Returns:
        dict chứa slide_count, shapes, text_content
    """
    if Presentation is None:
        print("Lỗi: Cần cài đặt python-pptx")
        print("Chạy: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(pptx_path)
    metadata = {
        "file": Path(pptx_path).name,
        "slide_count": len(prs.slides),
        "shapes": [],
        "text_content": []
    }

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            shape_data = {
                "slide": slide_idx,
                "type": shape.shape_type.name if hasattr(shape.shape_type, "name") else str(shape.shape_type),
                "left": round(shape.left / Inches(1), 3),
                "top": round(shape.top / Inches(1), 3),
                "width": round(shape.width / Inches(1), 3),
                "height": round(shape.height / Inches(1), 3),
            }

            # Trích xuất màu nếu có
            try:
                if hasattr(shape, "fill") and shape.fill.type is not None:
                    if shape.fill.fore_color and shape.fill.fore_color.rgb:
                        shape_data["fill_color"] = str(shape.fill.fore_color.rgb)
            except Exception:
                pass

            # Trích xuất text nếu có
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shape_data["text"] = text
                    metadata["text_content"].append(text)

            metadata["shapes"].append(shape_data)

    return metadata


def extract_excel_metadata(xlsx_path: str) -> dict:
    """
    Trích xuất metadata từ Excel.

    Returns:
        dict chứa sheet_names, sheets với headers và row_count
    """
    if load_workbook is None:
        print("Lỗi: Cần cài đặt openpyxl")
        print("Chạy: pip install openpyxl")
        sys.exit(1)

    wb = load_workbook(xlsx_path, read_only=True)
    metadata = {
        "file": Path(xlsx_path).name,
        "sheet_names": wb.sheetnames,
        "sheets": {}
    }

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        # Lấy headers từ dòng đầu
        headers = []
        for cell in next(ws.iter_rows(min_row=1, max_row=1), []):
            if cell.value:
                headers.append(str(cell.value))

        # Đếm số dòng có dữ liệu
        row_count = 0
        for row in ws.iter_rows(min_row=2):
            if any(cell.value for cell in row):
                row_count += 1

        metadata["sheets"][sheet_name] = {
            "headers": headers,
            "row_count": row_count,
            "col_count": len(headers)
        }

    wb.close()
    return metadata


def main():
    parser = argparse.ArgumentParser(
        description="Trích xuất metadata từ PPTX/Excel"
    )
    parser.add_argument(
        "file",
        help="Đường dẫn file PPTX hoặc Excel"
    )
    parser.add_argument(
        "-o", "--output",
        help="File output JSON (mặc định: stdout)"
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        help="Format JSON đẹp"
    )

    args = parser.parse_args()

    file_path = Path(args.file)
    if not file_path.exists():
        print(f"Lỗi: Không tìm thấy file {file_path}")
        sys.exit(1)

    suffix = file_path.suffix.lower()

    if suffix == ".pptx":
        metadata = extract_pptx_metadata(str(file_path))
    elif suffix in (".xlsx", ".xls"):
        metadata = extract_excel_metadata(str(file_path))
    else:
        print(f"Lỗi: Định dạng không hỗ trợ: {suffix}")
        print("Hỗ trợ: .pptx, .xlsx")
        sys.exit(1)

    # Output
    indent = 2 if args.pretty else None
    json_output = json.dumps(metadata, indent=indent, ensure_ascii=False)

    if args.output:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(json_output)
        print(f"Đã lưu metadata vào: {output_path}")
    else:
        print(json_output)


if __name__ == "__main__":
    main()
